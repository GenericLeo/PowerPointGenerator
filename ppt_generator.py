"""
PowerPoint Generator for Image Groups
Creates presentation slides from grouped images with optimized layout
"""

import os
from typing import List, Dict, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image as PILImage
import math


class PowerPointGenerator:
    """Generates PowerPoint presentations from grouped images"""
    
    def __init__(self, image_index):
        """
        Initialize PowerPoint generator
        
        Args:
            image_index: ImageIndex instance with indexed images
        """
        self.image_index = image_index
        
        # Custom layout profiles
        self.layout_profiles = {}
        
        # Standard slide dimensions (16:9)
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        
        # Margins and spacing
        self.margin_top = Inches(0.75)
        self.margin_bottom = Inches(0.5)
        self.margin_left = Inches(0.5)
        self.margin_right = Inches(0.5)
        
        # Text height for filename labels
        self.label_height = Inches(0.4)
        self.label_spacing = Inches(0.1)
        
        # Available area for images
        self.available_width = self.slide_width - self.margin_left - self.margin_right
        self.available_height = self.slide_height - self.margin_top - self.margin_bottom
    
    def calculate_grid_layout(self, num_images: int) -> Tuple[int, int]:
        """
        Calculate optimal grid layout (rows, cols) for given number of images
        
        Args:
            num_images: Number of images to fit on slide
            
        Returns:
            Tuple of (rows, cols)
        """
        if num_images == 0:
            return (0, 0)
        elif num_images == 1:
            return (1, 1)
        elif num_images == 2:
            return (1, 2)
        elif num_images <= 4:
            return (2, 2)
        elif num_images <= 6:
            return (2, 3)
        elif num_images <= 9:
            return (3, 3)
        elif num_images <= 12:
            return (3, 4)
        elif num_images <= 16:
            return (4, 4)
        elif num_images <= 20:
            return (4, 5)
        else:
            # For very large groups, use square-ish grid
            cols = math.ceil(math.sqrt(num_images))
            rows = math.ceil(num_images / cols)
            return (rows, cols)
    
    def calculate_image_dimensions(self, rows: int, cols: int, 
                                   image_aspect_ratio: float) -> Tuple[float, float]:
        """
        Calculate optimal image dimensions to fit in grid cell
        
        Args:
            rows: Number of rows in grid
            cols: Number of columns in grid
            image_aspect_ratio: Width/height ratio of image
            
        Returns:
            Tuple of (width, height) in Inches
        """
        # Calculate cell dimensions
        cell_width = self.available_width / cols
        cell_height = (self.available_height - (self.label_height + self.label_spacing) * rows) / rows
        
        # Calculate image size that fits in cell while maintaining aspect ratio
        if cell_width / cell_height > image_aspect_ratio:
            # Height constrained
            img_height = cell_height - self.label_height - self.label_spacing
            img_width = img_height * image_aspect_ratio
        else:
            # Width constrained
            img_width = cell_width * 0.95  # 95% to leave small gap
            img_height = img_width / image_aspect_ratio
            
            # Ensure height fits
            max_img_height = cell_height - self.label_height - self.label_spacing
            if img_height > max_img_height:
                img_height = max_img_height
                img_width = img_height * image_aspect_ratio
        
        return (img_width, img_height)
    
    def add_group_slide(self, prs: Presentation, group_label: str, images: List[Dict]):
        """
        Add a slide for a group of images with optimized layout
        
        Args:
            prs: Presentation object
            group_label: Label for the group (e.g., "0001", "MAP1", "SPEC1")
            images: List of image dictionaries from index
        """
        if not images:
            return
        
        # Check if there's a custom layout profile for this group
        if group_label in self.layout_profiles:
            profile = self.layout_profiles[group_label]
            layout_type = profile.get('type', 'auto')
            
            if layout_type == 'visual':
                # Use visual custom layout
                self.add_visual_custom_slide(prs, group_label, images, profile)
                return
            elif layout_type == 'grid':
                self.add_standard_slide(prs, group_label, images)
                return
            elif layout_type == 'horizontal':
                self.add_spectrum_slide(prs, group_label, images)
                return
            elif layout_type == 'mixed':
                # Force mixed layout
                from image_uploader import ImageIdentifier
                spectrum_images = [img for img in images 
                                  if img.get('metadata', {}).get('identifier', '') == 'Spectrum' or 
                                  img.get('metadata', {}).get('identifier', '') in ImageIdentifier.SPECTRUM_IDENTIFIERS]
                other_images = [img for img in images 
                               if img.get('metadata', {}).get('identifier', '') != 'Spectrum' and
                               img.get('metadata', {}).get('identifier', '') not in ImageIdentifier.SPECTRUM_IDENTIFIERS]
                
                if not spectrum_images:
                    spectrum_images = images[:3]  # Treat first 3 as spectrum-like
                    other_images = images[3:]
                if not other_images and len(images) > 3:
                    other_images = images[-2:]  # Treat last 2 as other
                    spectrum_images = images[:-2]
                
                self.add_mixed_slide(prs, group_label, spectrum_images if spectrum_images else [], 
                                   other_images if other_images else [])
                return
            # 'custom' and 'auto' fall through to default behavior
        
        # Default auto-detection behavior
        from image_uploader import ImageIdentifier
        
        spectrum_images = [img for img in images 
                          if img.get('metadata', {}).get('identifier', '') == 'Spectrum' or
                          img.get('metadata', {}).get('identifier', '') in ImageIdentifier.SPECTRUM_IDENTIFIERS]
        other_images = [img for img in images 
                       if img.get('metadata', {}).get('identifier', '') != 'Spectrum' and
                       img.get('metadata', {}).get('identifier', '') not in ImageIdentifier.SPECTRUM_IDENTIFIERS]
        
        if spectrum_images and other_images:
            # Mixed group: use special layout
            self.add_mixed_slide(prs, group_label, spectrum_images, other_images)
        elif spectrum_images:
            # Spectrum only: use horizontal layout
            self.add_spectrum_slide(prs, group_label, spectrum_images)
        else:
            # Standard images: use grid layout
            self.add_standard_slide(prs, group_label, images)
    
    def add_visual_custom_slide(self, prs: Presentation, group_label: str, 
                               images: List[Dict], profile: Dict):
        """
        Add a slide with visual custom layout defined by user-drawn regions
        
        Args:
            prs: Presentation object
            group_label: Label for the group
            images: List of all image dictionaries for this group
            profile: Layout profile containing regions configuration
        """
        from image_uploader import ImageIdentifier
        
        # Get regions from profile
        regions = profile.get('regions', [])
        if not regions:
            # Fall back to standard layout if no regions defined
            self.add_standard_slide(prs, group_label, images)
            return
        
        # Group images by their identifier
        images_by_identifier = {}
        for img in images:
            identifier = img.get('metadata', {}).get('identifier', 'Unknown')
            if identifier not in images_by_identifier:
                images_by_identifier[identifier] = []
            images_by_identifier[identifier].append(img)
        
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            self.margin_left,
            Inches(0.25),
            self.available_width,
            Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"Group {group_label}"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Process each region
        for region in regions:
            region_identifier = region['identifier']
            region_images = images_by_identifier.get(region_identifier, [])
            
            if not region_images:
                continue
            
            # Calculate region position and size (regions are stored as percentages 0-1)
            region_x = self.margin_left + region['x1'] * self.available_width
            region_y = self.margin_top + region['y1'] * self.available_height
            region_width = (region['x2'] - region['x1']) * self.available_width
            region_height = (region['y2'] - region['y1']) * self.available_height
            
            # Calculate grid for images in this region
            num_region_images = len(region_images)
            rows, cols = self.calculate_grid_layout(num_region_images)
            
            # Calculate cell dimensions within region
            cell_width = region_width / cols
            cell_height = region_height / rows
            
            # Place images in region grid
            for idx, img in enumerate(region_images):
                row = idx // cols
                col = idx % cols
                
                # Get image aspect ratio
                img_aspect_ratio = img['width'] / img['height'] if img['height'] > 0 else 1.0
                
                # Calculate available space (leaving room for label)
                available_cell_height = cell_height - self.label_height - self.label_spacing
                
                # Calculate image dimensions maintaining aspect ratio
                if cell_width / available_cell_height > img_aspect_ratio:
                    # Height constrained
                    img_height = available_cell_height
                    img_width = img_height * img_aspect_ratio
                else:
                    # Width constrained
                    img_width = cell_width * 0.95
                    img_height = img_width / img_aspect_ratio
                    if img_height > available_cell_height:
                        img_height = available_cell_height
                        img_width = img_height * img_aspect_ratio
                
                # Calculate position (centered in cell)
                x = region_x + col * cell_width + (cell_width - img_width) / 2
                y = region_y + row * cell_height + (available_cell_height - img_height) / 2
                
                # Add image
                try:
                    slide.shapes.add_picture(
                        img['filepath'],
                        x, y,
                        width=img_width,
                        height=img_height
                    )
                    
                    # Add filename label below image
                    label_y = y + img_height + self.label_spacing
                    label_box = slide.shapes.add_textbox(
                        x,
                        label_y,
                        img_width,
                        self.label_height
                    )
                    label_frame = label_box.text_frame
                    label_frame.text = img['filename']
                    label_frame.paragraphs[0].font.size = Pt(8)
                    label_frame.word_wrap = True
                    label_frame.paragraphs[0].alignment = 1  # Center alignment
                    
                except Exception as e:
                    print(f"Error adding image {img['filename']} to region: {e}")
    
    def add_mixed_slide(self, prs: Presentation, group_label: str, 
                       spectrum_images: List[Dict], other_images: List[Dict]):
        """
        Add a slide with mixed content: spectra on right 2/3, up to 2 other images on left 1/3
        
        Args:
            prs: Presentation object
            group_label: Label for the group
            spectrum_images: List of spectrum image dictionaries (max 3 per slide)
            other_images: List of other image dictionaries (max 2 per slide)
        """
        # Process in batches: max 3 spectra and 2 other images per slide
        spectrum_batch_size = 3
        other_batch_size = 2
        
        max_slides = max(
            (len(spectrum_images) + spectrum_batch_size - 1) // spectrum_batch_size,
            (len(other_images) + other_batch_size - 1) // other_batch_size
        )
        
        for slide_num in range(max_slides):
            # Get batches for this slide
            spec_start = slide_num * spectrum_batch_size
            spec_batch = spectrum_images[spec_start:spec_start + spectrum_batch_size]
            
            other_start = slide_num * other_batch_size
            other_batch = other_images[other_start:other_start + other_batch_size]
            
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_text = f"Group {group_label}"
            if max_slides > 1:
                title_text += f" (Part {slide_num + 1}/{max_slides})"
            
            title_box = slide.shapes.add_textbox(
                self.margin_left,
                Inches(0.25),
                self.available_width,
                Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # Calculate layout areas
            left_width = self.available_width / 3  # Left 1/3 for other images
            right_width = (self.available_width * 2) / 3  # Right 2/3 for spectra
            
            # Add other images on left (max 2, stacked vertically)
            if other_batch:
                num_other = len(other_batch)
                other_slice_height = self.available_height / num_other
                
                for idx, img in enumerate(other_batch):
                    y_pos = self.margin_top + idx * other_slice_height
                    img_aspect_ratio = img['width'] / img['height'] if img['height'] > 0 else 1.0
                    
                    # Reserve space for label
                    available_slice_height = other_slice_height - self.label_height - self.label_spacing
                    
                    # Calculate dimensions to fit in left 1/3
                    img_width = left_width * 0.9
                    img_height = img_width / img_aspect_ratio
                    
                    if img_height > available_slice_height:
                        img_height = available_slice_height
                        img_width = img_height * img_aspect_ratio
                    
                    # Center in left area
                    x_pos = self.margin_left + (left_width - img_width) / 2
                    
                    try:
                        slide.shapes.add_picture(
                            img['filepath'],
                            x_pos, y_pos,
                            width=img_width,
                            height=img_height
                        )
                        
                        # Add label
                        label_y = y_pos + img_height + self.label_spacing
                        label_box = slide.shapes.add_textbox(
                            self.margin_left,
                            label_y,
                            left_width,
                            self.label_height
                        )
                        label_frame = label_box.text_frame
                        label_frame.text = img['filename']
                        label_frame.paragraphs[0].font.size = Pt(8)
                        label_frame.word_wrap = True
                        label_frame.paragraphs[0].alignment = 1  # Center
                    except Exception as e:
                        print(f"Error adding image {img['filename']}: {e}")
            
            # Add spectrum images on right (max 3, horizontal slices)
            if spec_batch:
                num_spectra = len(spec_batch)
                spec_slice_height = self.available_height / num_spectra
                right_x_start = self.margin_left + left_width
                
                for idx, img in enumerate(spec_batch):
                    y_pos = self.margin_top + idx * spec_slice_height
                    img_aspect_ratio = img['width'] / img['height'] if img['height'] > 0 else 4.0
                    
                    # Reserve space for label
                    available_slice_height = spec_slice_height - self.label_height - self.label_spacing
                    
                    # Calculate dimensions to fit in right 2/3
                    img_width = right_width * 0.95
                    img_height = img_width / img_aspect_ratio
                    
                    if img_height > available_slice_height:
                        img_height = available_slice_height
                        img_width = img_height * img_aspect_ratio
                    
                    # Center in right area
                    x_pos = right_x_start + (right_width - img_width) / 2
                    
                    try:
                        slide.shapes.add_picture(
                            img['filepath'],
                            x_pos, y_pos,
                            width=img_width,
                            height=img_height
                        )
                        
                        # Add label
                        label_y = y_pos + img_height + self.label_spacing
                        label_box = slide.shapes.add_textbox(
                            right_x_start,
                            label_y,
                            right_width,
                            self.label_height
                        )
                        label_frame = label_box.text_frame
                        label_frame.text = img['filename']
                        label_frame.paragraphs[0].font.size = Pt(10)
                        label_frame.word_wrap = True
                        label_frame.paragraphs[0].alignment = 1  # Center
                    except Exception as e:
                        print(f"Error adding spectrum {img['filename']}: {e}")
    
    def add_spectrum_slide(self, prs: Presentation, group_label: str, images: List[Dict]):
        """
        Add a slide for spectrum images with horizontal layout (max 3 per slide)
        
        Args:
            prs: Presentation object
            group_label: Label for the group (e.g., "SPEC1")
            images: List of spectrum image dictionaries
        """
        if not images:
            return
        
        # Process in batches of 3
        for batch_start in range(0, len(images), 3):
            batch = images[batch_start:batch_start + 3]
            num_spectra = len(batch)
            
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_text = f"Group {group_label}"
            if len(images) > 3:
                batch_num = (batch_start // 3) + 1
                total_batches = (len(images) + 2) // 3
                title_text += f" (Part {batch_num}/{total_batches})"
            
            title_box = slide.shapes.add_textbox(
                self.margin_left,
                Inches(0.25),
                self.available_width,
                Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # Calculate horizontal slice layout
            slice_height = self.available_height / num_spectra
            
            for idx, img in enumerate(batch):
                # Calculate slice position
                y_pos = self.margin_top + idx * slice_height
                
                # Get image aspect ratio
                img_aspect_ratio = img['width'] / img['height'] if img['height'] > 0 else 4.0
                
                # Reserve space for filename label at bottom of slice
                available_slice_height = slice_height - self.label_height - self.label_spacing
                
                # Calculate image dimensions to fit in horizontal slice
                # Try to use full width first
                img_width = self.available_width * 0.95  # 95% of available width
                img_height = img_width / img_aspect_ratio
                
                # Check if height fits in slice
                if img_height > available_slice_height:
                    img_height = available_slice_height
                    img_width = img_height * img_aspect_ratio
                
                # Center horizontally
                x_pos = self.margin_left + (self.available_width - img_width) / 2
                
                # Add image
                try:
                    slide.shapes.add_picture(
                        img['filepath'],
                        x_pos, y_pos,
                        width=img_width,
                        height=img_height
                    )
                    
                    # Add filename label below image
                    label_y = y_pos + img_height + self.label_spacing
                    label_box = slide.shapes.add_textbox(
                        self.margin_left,
                        label_y,
                        self.available_width,
                        self.label_height
                    )
                    label_frame = label_box.text_frame
                    label_frame.text = img['filename']
                    label_frame.paragraphs[0].font.size = Pt(10)
                    label_frame.word_wrap = True
                    
                    # Center align text
                    label_frame.paragraphs[0].alignment = 1  # Center alignment
                    
                except Exception as e:
                    print(f"Error adding spectrum image {img['filename']}: {e}")
    
    def add_standard_slide(self, prs: Presentation, group_label: str, images: List[Dict]):
        """
        Add a slide for standard images with grid layout
        
        Args:
            prs: Presentation object
            group_label: Label for the group (e.g., "0001", "MAP1")
            images: List of image dictionaries from index
        """
        if not images:
            return
        
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            self.margin_left,
            Inches(0.25),
            self.available_width,
            Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"Group {group_label}"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Calculate grid layout
        num_images = len(images)
        rows, cols = self.calculate_grid_layout(num_images)
        
        # Get average aspect ratio of images for layout calculation
        aspect_ratios = []
        for img in images:
            if img['width'] and img['height']:
                aspect_ratios.append(img['width'] / img['height'])
        
        avg_aspect_ratio = sum(aspect_ratios) / len(aspect_ratios) if aspect_ratios else 1.0
        
        # Calculate cell dimensions
        cell_width = self.available_width / cols
        cell_height = self.available_height / rows
        
        # Place images in grid
        for idx, img in enumerate(images):
            row = idx // cols
            col = idx % cols
            
            # Get image aspect ratio
            img_aspect_ratio = img['width'] / img['height'] if img['height'] > 0 else 1.0
            
            # Calculate image dimensions
            img_width, img_height = self.calculate_image_dimensions(rows, cols, img_aspect_ratio)
            
            # Calculate position (centered in cell)
            x = self.margin_left + col * cell_width + (cell_width - img_width) / 2
            y = self.margin_top + row * cell_height + (cell_height - img_height - self.label_height - self.label_spacing) / 2
            
            # Add image
            try:
                slide.shapes.add_picture(
                    img['filepath'],
                    x, y,
                    width=img_width,
                    height=img_height
                )
                
                # Add filename label below image
                label_y = y + img_height + self.label_spacing
                label_box = slide.shapes.add_textbox(
                    x,
                    label_y,
                    img_width,
                    self.label_height
                )
                label_frame = label_box.text_frame
                label_frame.text = img['filename']
                label_frame.paragraphs[0].font.size = Pt(8)
                label_frame.word_wrap = True
                
                # Center align text
                label_frame.paragraphs[0].alignment = 1  # Center alignment
                
            except Exception as e:
                print(f"Error adding image {img['filename']}: {e}")
    
    def generate_presentation(self, output_path: str, groups_to_include: List[str] = None) -> bool:
        """
        Generate PowerPoint presentation with one slide per group
        
        Args:
            output_path: Path to save the presentation
            groups_to_include: Optional list of group labels to include (e.g., ["0001", "MAP1"])
                              If None, includes all groups
            
        Returns:
            True if successful, False otherwise
        """
        try:
            from image_uploader import ImageIdentifier
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = self.slide_width
            prs.slide_height = self.slide_height
            
            # Get all images grouped
            images = self.image_index.images
            
            # Group images by their formatted label
            groups_dict = {}
            for img in images:
                metadata = img.get('metadata', {})
                numerical_prefix = metadata.get('numerical_prefix')
                identifier = metadata.get('identifier')
                
                if not numerical_prefix:
                    continue  # Skip ungrouped images
                
                # Format group label
                group_label = ImageIdentifier.format_group_label(numerical_prefix, identifier)
                
                # Filter by groups_to_include if specified
                if groups_to_include and group_label not in groups_to_include:
                    continue
                
                if group_label not in groups_dict:
                    groups_dict[group_label] = []
                groups_dict[group_label].append(img)
            
            # Sort groups
            sorted_groups = sorted(groups_dict.keys())
            
            # Add a slide for each group
            for group_label in sorted_groups:
                group_images = groups_dict[group_label]
                self.add_group_slide(prs, group_label, group_images)
            
            # Save presentation
            prs.save(output_path)
            print(f"PowerPoint saved to: {output_path}")
            print(f"Total slides: {len(sorted_groups)}")
            
            return True
            
        except Exception as e:
            print(f"Error generating presentation: {e}")
            import traceback
            traceback.print_exc()
            return False
